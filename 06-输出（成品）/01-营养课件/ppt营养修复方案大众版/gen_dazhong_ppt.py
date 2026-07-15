#!/usr/bin/env python3
"""基于大众版最新讲稿生成PPT——字体适中、无动画、直接PPT格式"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色 ──
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT  = RGBColor(0x00, 0x6D, 0xAA)
ACCENT2 = RGBColor(0xF0, 0x8C, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF5, 0xF7, 0xFA)
GRAY    = RGBColor(0x66, 0x66, 0x66)
GREEN   = RGBColor(0x27, 0xAE, 0x60)
RED     = RGBColor(0xE7, 0x4C, 0x3C)
BG_DARK  = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG  = RGBColor(0xF8, 0xFA, 0xFC)
BLUE_BG  = RGBColor(0xEB, 0xF5, 0xFF)
ORANGE_BG = RGBColor(0xFD, 0xF0, 0xE0)
GREEN_BG  = RGBColor(0xE8, 0xF8, 0xF0)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── 字体大小常量 ──
TITLE    = 36   # 页面标题
SUBTITLE = 22   # 副标题
BODY     = 20   # 正文
SMALL    = 18   # 小字/表格
CAPTION  = 16   # 注释
LARGE    = 28   # 大号强调
HUGE     = 36   # 巨型数字/封面主标题
COVER_TITLE = 52  # 封面标题

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(2)
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=BODY, bold=False,
                color=DARK, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, default_size=BODY,
                  default_color=DARK, default_bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, size, bold, color = item, default_size, default_bold, default_color
        else:
            text = item.get('text', '')
            size = item.get('size', default_size)
            bold = item.get('bold', default_bold)
            color = item.get('color', default_color)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_before = Pt(item.get('space', 4) if isinstance(item, dict) else 4)
    return txBox

def add_page_title(slide, title_text, subtitle_text=None):
    add_textbox(slide, Inches(0.6), Inches(0.3), Inches(11), Inches(0.8),
                title_text, font_size=TITLE, bold=True, color=DARK)
    add_shape(slide, Inches(0.6), Inches(1.0), Inches(1.5), Inches(0.04), ACCENT)
    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(1.15), Inches(11), Inches(0.5),
                    subtitle_text, font_size=SUBTITLE, color=GRAY)

# ══════════════════════════════════════════════════════════════
# P1 · 封面
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.25), H, ACCENT)
add_textbox(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(1.5),
            '人为什么会生病', font_size=COVER_TITLE, bold=True, color=WHITE)
add_textbox(slide, Inches(1.2), Inches(3.2), Inches(10), Inches(0.8),
            '从隐性饥饿到细胞修复', font_size=HUGE, color=RGBColor(0x8E, 0xC6, 0xFF))
add_textbox(slide, Inches(1.2), Inches(5.8), Inches(8), Inches(0.6),
            '给你的细胞它需要的材料，它会自己修好自己', font_size=BODY, color=GRAY)
add_shape(slide, Inches(1.2), Inches(5.5), Inches(2.5), Inches(0.05), ACCENT)

# ══════════════════════════════════════════════════════════════
# P2 · 扎心事实
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '每年体检的扎心真相')

items = [
    ('📋 体检报告', '上上下下的箭头\n每年都有新问题', ACCENT, BLUE_BG),
    ('👨‍⚕️ 医生怎么说', '"没什么大问题\n注意饮食多运动"', GRAY, CARD_BG),
    ('❓ 然后呢？', '到底怎么注意？\n没人告诉你', RED, RGBColor(0xFD, 0xE0, 0xE0)),
]
for i, (title, desc, color, bg) in enumerate(items):
    x = Inches(0.6 + i * 4.2)
    add_rounded_rect(slide, x, Inches(1.8), Inches(3.8), Inches(3.0), bg)
    add_shape(slide, x, Inches(1.8), Inches(3.8), Inches(0.06), color)
    add_textbox(slide, x + Inches(0.2), Inches(2.0), Inches(3.4), Inches(0.5),
                title, font_size=LARGE, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(2.6), Inches(3.4), Inches(2.0),
                desc, font_size=BODY, color=DARK, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(0.8),
            '你不是一个人这样——是整个体系里缺了一环：没人从底层告诉你身体到底需要什么',
            font_size=SUBTITLE, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# P3 · 今晚目标
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '今晚的目标地图')

goals = [
    ('❶', '身体到底缺了什么', '为什么每天吃三顿饭，细胞还在"喊饿"', ACCENT),
    ('❷', '应该怎么补', '一个简单的基础方案，先把最缺的补上', ACCENT2),
    ('❸', '怎么省钱吃', '还可能赚钱——长客会计划', GREEN),
]
for i, (num, title, desc, color) in enumerate(goals):
    y = Inches(1.8 + i * 1.7)
    add_rounded_rect(slide, Inches(0.6), y, Inches(12), Inches(1.3), CARD_BG)
    add_shape(slide, Inches(0.6), y, Inches(0.06), Inches(1.3), color)
    add_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(0.8), Inches(0.8),
                num, font_size=HUGE, bold=True, color=color)
    add_textbox(slide, Inches(2.0), y + Inches(0.1), Inches(4.5), Inches(0.5),
                title, font_size=LARGE, bold=True, color=DARK)
    add_textbox(slide, Inches(2.0), y + Inches(0.65), Inches(10), Inches(0.5),
                desc, font_size=BODY, color=GRAY)

# ══════════════════════════════════════════════════════════════
# P4 · 从细胞说起
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '从细胞说起', '细胞健康 = 人健康')

steps_list = ['细胞', '组织', '器官', '系统', '人']
for i, s in enumerate(steps_list):
    x = Inches(0.8 + i * 2.5)
    add_rounded_rect(slide, x, Inches(1.8), Inches(1.8), Inches(1.0), ACCENT)
    add_textbox(slide, x, Inches(1.9), Inches(1.8), Inches(0.8),
                s, font_size=LARGE, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < len(steps_list) - 1:
        add_textbox(slide, x + Inches(1.8), Inches(1.9), Inches(0.7), Inches(0.8),
                    '→', font_size=HUGE, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(2.5), Inches(3.3), Inches(8), Inches(1.5), BLUE_BG)
add_shape(slide, Inches(2.5), Inches(3.3), Inches(8), Inches(0.06), ACCENT)
lines_cell = [
    {'text': '🧬 核心公式：细胞健康 = 人健康', 'size': LARGE, 'bold': True, 'color': ACCENT, 'space': 8},
    {'text': '如果每个细胞都是健康的 → 组织、器官、人就是健康的', 'size': BODY, 'color': DARK, 'space': 6},
]
add_multiline(slide, Inches(3), Inches(3.5), Inches(7), Inches(1.2), lines_cell)

add_textbox(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(0.6),
            '细胞出了问题，吃再多药、做再多检查，都只是治标不治本',
            font_size=BODY, color=GRAY, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# P5 · 7大营养素
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '细胞需要什么？——7大营养素')

nuts = [
    ('💧', '水', '身体的溶剂·所有化学反应的介质', ACCENT),
    ('🧱', '蛋白质', '细胞修复的砖·免疫球蛋白·肌肉', ACCENT),
    ('🔧', '维生素', '代谢辅酶·催化所有反应', ACCENT2),
    ('⛏️', '矿物质', '结构组成·酶活化', ACCENT2),
    ('🔥', '碳水', '燃料·提供能量', GREEN),
    ('🛢️', '脂肪', '细胞膜结构·储能与保温', GREEN),
    ('🧹', '膳食纤维', '肠道清道夫·排出毒素', GRAY),
]
for i, (emoji, name, desc, color) in enumerate(nuts):
    col = i % 4
    row = i // 4
    x = Inches(0.6 + col * 3.15)
    y = Inches(1.6 + row * 2.7)
    add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.3), CARD_BG)
    add_shape(slide, x, y, Inches(2.8), Inches(0.05), color)
    add_textbox(slide, x, y + Inches(0.15), Inches(2.8), Inches(0.5),
                f'{emoji}  {name}', font_size=LARGE, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.8), Inches(2.5), Inches(1.3),
                desc, font_size=SMALL, color=DARK, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# P6 · 一句话逻辑
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.25), H, ACCENT)

add_textbox(slide, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
            '健康管理的底层逻辑', font_size=TITLE, bold=True, color=WHITE)
add_shape(slide, Inches(1.2), Inches(1.5), Inches(1.5), Inches(0.04), ACCENT)

add_rounded_rect(slide, Inches(1.2), Inches(2.2), Inches(11), Inches(2.2), RGBColor(0x15, 0x25, 0x40))
add_shape(slide, Inches(1.2), Inches(2.2), Inches(11), Inches(0.06), ACCENT2)
add_textbox(slide, Inches(1.8), Inches(2.5), Inches(9.5), Inches(1.6),
            '原材料不够，细胞就没法正常干活。\n人就会出问题。',
            font_size=HUGE, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

lines_logic = [
    {'text': '你的身体每分每秒都在用营养素做修复、做代谢、做免疫', 'size': BODY, 'color': RGBColor(0xAA, 0xBB, 0xDD)},
    {'text': '原材料够了，身体自己会修自己', 'size': BODY, 'bold': True, 'color': RGBColor(0xCC, 0xDD, 0xFF)},
    {'text': '原材料不够，修不好，问题就一点点积累', 'size': BODY, 'color': RGBColor(0xAA, 0xBB, 0xDD)},
]
add_multiline(slide, Inches(1.2), Inches(4.8), Inches(11), Inches(2.0), lines_logic)

# ══════════════════════════════════════════════════════════════
# P7 · 三重剥削
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '吃 ≠ 吃饱——三重剥削', '你吃进去的 ≠ 细胞收到的')

levels = [
    ('第一重：农业空心化', '−20%', '大棚速生·蛋白质缩水15-20%', ACCENT),
    ('第二重：烹饪毒化', '−30%', '外卖 AGEs 毒化报废', ACCENT2),
    ('第三重：消化无能', '−20%', '高压→胃酸不足→大颗粒流失', RED),
]
for i, (title, pct, desc, color) in enumerate(levels):
    x = Inches(0.6 + i * 4.2)
    y = Inches(1.6)
    add_rounded_rect(slide, x, y, Inches(3.8), Inches(3.5), CARD_BG)
    add_shape(slide, x, y, Inches(3.8), Inches(0.06), color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.5),
                title, font_size=BODY, bold=True, color=color)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(0.9),
                pct, font_size=48, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(2.0), Inches(3.4), Inches(1.2),
                desc, font_size=BODY, color=DARK, alignment=PP_ALIGN.CENTER)

add_rounded_rect(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(0.7), RGBColor(0xFD, 0xE0, 0xE0))
add_textbox(slide, Inches(1.0), Inches(5.55), Inches(11), Inches(0.6),
            '⚡ 三关下来，到细胞的可能不到三成——你的胃吃饱了，但你的细胞还饿着',
            font_size=BODY, bold=True, color=RED, alignment=PP_ALIGN.CENTER)

# + 营养素缺口表
headers = ['营养素', '标称', '实到', '修复需要', '缺口']
rows_p7 = [
    ['蛋白质',   '55g',    '25-27g',  '72-96g',    '🔴缺45-70g'],
    ['Omega-3',  '≈0.1g',  '≈0g',     '3.0-4.0g',  '🔴缺3g+'],
    ['镁',       '150mg',  '40-50mg', '500-800mg', '🔴缺450-750mg'],
    ['维C',      '≈20mg',  '<10mg',   '500-1,000mg','🔴缺490-990mg'],
]
col_w = [Inches(1.6), Inches(1.3), Inches(1.3), Inches(1.8), Inches(2.2)]
col_s = [Inches(0.8)]
for w in col_w[:-1]:
    col_s.append(col_s[-1] + w)
y_t = Inches(6.3)
add_shape(slide, Inches(0.8), y_t, Inches(8.2), Inches(0.55), ACCENT)
for hdr, cs, cw in zip(headers, col_s, col_w):
    add_textbox(slide, cs, y_t + Inches(0.03), cw, Inches(0.5),
                hdr, font_size=CAPTION, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
for r_idx, row in enumerate(rows_p7):
    y_r = y_t + Inches(0.55) + r_idx * Inches(0.55)
    bg_c = RGBColor(0xF8, 0xFA, 0xFC) if r_idx % 2 == 0 else WHITE
    add_shape(slide, Inches(0.8), y_r, Inches(8.2), Inches(0.55), bg_c)
    for c_idx, (cell, cs, cw) in enumerate(zip(row, col_s, col_w)):
        c = RED if '🔴' in cell else DARK
        add_textbox(slide, cs, y_r + Inches(0.03), cw, Inches(0.5),
                    cell, font_size=CAPTION, bold=(c_idx==0), color=c, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# P8 · 隐性饥饿
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '隐性饥饿——你的细胞在挨饿')

add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(5.3), BLUE_BG)
add_shape(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.06), ACCENT)
lines_hunger = [
    {'text': '什么是隐性饥饿？', 'size': LARGE, 'bold': True, 'color': ACCENT, 'space': 8},
    {'text': '世界卫生组织定义：微量营养素（维生素和矿物质）摄入不足，但你完全感觉不到。', 'size': BODY, 'color': DARK, 'space': 10},
    {'text': '', 'size': 8, 'color': DARK, 'space': 2},
    {'text': '它不疼不痒——不像没吃饭那样肚子饿得咕咕叫。', 'size': BODY, 'color': DARK, 'space': 6},
    {'text': '每天缺一点点，一年两年，问题就出来了。', 'size': BODY, 'bold': True, 'color': RED, 'space': 10},
    {'text': '', 'size': 8, 'color': DARK, 'space': 2},
    {'text': '就像水蛭吸血——每天吸，天天吸，等发现的时候已经出问题了。', 'size': BODY, 'color': DARK, 'space': 6},
    {'text': '慢性病就是这样来的：不是突然发生的，是突然发现的。', 'size': LARGE, 'bold': True, 'color': ACCENT, 'space': 12},
]
add_multiline(slide, Inches(1.0), Inches(1.8), Inches(11), Inches(4.5), lines_hunger)

# ══════════════════════════════════════════════════════════════
# P9 · 为什么感觉不到
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.25), H, ACCENT)

add_textbox(slide, Inches(1.2), Inches(0.8), Inches(10), Inches(0.8),
            '为什么你感觉不到？', font_size=TITLE, bold=True, color=WHITE)
add_shape(slide, Inches(1.2), Inches(1.5), Inches(1.5), Inches(0.04), ACCENT)

add_rounded_rect(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(1.6), RGBColor(0x15, 0x25, 0x40))
add_shape(slide, Inches(1.2), Inches(2.0), Inches(11), Inches(0.06), ACCENT2)
add_textbox(slide, Inches(1.8), Inches(2.2), Inches(9.5), Inches(1.2),
            '慢性病不是突然发生的，是突然发现的。',
            font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

examples_p9 = [
    ('🩺 高血压', '血管壁慢慢变硬了十年\n昨天一量才发现'),
    ('🩸 糖尿病', '胰岛细胞累趴了好几年\n上个月体检才发现'),
    ('🫁 脂肪肝', '肝细胞被脂肪包围了好几年\n这周B超才发现'),
]
for i, (title, desc) in enumerate(examples_p9):
    x = Inches(1.2 + i * 3.9)
    add_rounded_rect(slide, x, Inches(4.0), Inches(3.5), Inches(2.8), RGBColor(0x15, 0x25, 0x40))
    add_shape(slide, x, Inches(4.0), Inches(3.5), Inches(0.05), ACCENT if i == 0 else (ACCENT2 if i == 1 else RED))
    add_textbox(slide, x + Inches(0.3), Inches(4.2), Inches(2.9), Inches(0.5),
                title, font_size=LARGE, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    txBox = add_textbox(slide, x + Inches(0.3), Inches(4.9), Inches(2.9), Inches(1.6),
                        '', font_size=SMALL)
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, ll in enumerate(desc.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = ll
        p.font.size = Pt(BODY)
        p.font.color.rgb = RGBColor(0xBB, 0xCC, 0xEE)
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

# ══════════════════════════════════════════════════════════════
# P10 · 生命器官理论
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '身体的生存策略——生命器官 vs 非生命器官')

add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.0), GREEN_BG)
add_shape(slide, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.06), GREEN)
add_textbox(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.5),
            '✅ 生命器官（优先供给）', font_size=LARGE, bold=True, color=GREEN)
vital = ['❤️ 心脏', '🫁 肺', '🧠 大脑', '🫘 肝脏', '🧊 肾脏']
for i, v in enumerate(vital):
    add_textbox(slide, Inches(1.0), Inches(2.4 + i * 0.6), Inches(5), Inches(0.5),
                v, font_size=BODY, color=DARK)
add_textbox(slide, Inches(1.0), Inches(5.6), Inches(5), Inches(0.5),
            '少了任何一个，你当场就活不了', font_size=SMALL, color=GRAY)

add_rounded_rect(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.0), RGBColor(0xFD, 0xE0, 0xE0))
add_shape(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.06), RED)
add_textbox(slide, Inches(7.4), Inches(1.7), Inches(5), Inches(0.5),
            '⚠️ 非生命器官（被牺牲）', font_size=LARGE, bold=True, color=RED)
non_vital = ['🫃 胃', '🫀 乳腺', '🦋 甲状腺', '🧬 子宫/卵巢', '🦵 四肢', '🧴 皮肤']
for i, nv in enumerate(non_vital):
    add_textbox(slide, Inches(7.4), Inches(2.4 + i * 0.6), Inches(5), Inches(0.5),
                nv, font_size=BODY, color=DARK)
add_textbox(slide, Inches(7.4), Inches(5.6), Inches(5), Inches(0.5),
            '缺了不会马上死，但生活质量大幅下降', font_size=SMALL, color=GRAY)

# ══════════════════════════════════════════════════════════════
# P11 · 自检清单
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '你的身体在喊救命——自检清单')

symptoms = [
    ('口腔溃疡反复发作', '缺维生素B族', ACCENT),
    ('伤口愈合慢', '缺蛋白质', ACCENT),
    ('皮肤松弛、头发枯黄', '缺胶原蛋白原料', ACCENT2),
    ('乳腺增生', '营养被挪用了', ACCENT2),
    ('甲状腺结节', '营养被挪用了', RED),
    ('关节不舒服', '缺钙/镁/Omega-3', RED),
    ('总是觉得累', '缺B族/蛋白质/铁', GREEN),
    ('免疫力差容易感冒', '缺蛋白质/维C/锌', GREEN),
]
for i, (symp, reason, color) in enumerate(symptoms):
    col = i % 4
    row = i // 4
    x = Inches(0.6 + col * 3.15)
    y = Inches(1.6 + row * 2.5)
    add_rounded_rect(slide, x, y, Inches(2.8), Inches(2.1), CARD_BG)
    add_shape(slide, x, y, Inches(2.8), Inches(0.05), color)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.5), Inches(0.7),
                symp, font_size=BODY, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), y + Inches(1.1), Inches(2.5), Inches(0.7),
                reason, font_size=SMALL, color=color, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.5), ACCENT)
add_textbox(slide, Inches(1.0), Inches(6.82), Inches(11), Inches(0.45),
            '这些不是"亚健康"，是非生命器官在喊救命——原材料不够了',
            font_size=BODY, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# P12 · 算账
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '算一笔账——你每天缺多少？', '以60公斤成年人为例')

add_rounded_rect(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(4.5), BLUE_BG)
add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.06), ACCENT)
lines_water = [
    {'text': '💧 水', 'size': LARGE, 'bold': True, 'color': ACCENT, 'space': 6},
    {'text': '需要：1,800 ml/天  |  实际：1,000~1,200 ml/天', 'size': BODY, 'color': DARK, 'space': 10},
    {'text': '', 'size': 6, 'color': DARK, 'space': 2},
    {'text': '缺口：-600~800 ml/天', 'size': LARGE, 'bold': True, 'color': RED, 'space': 12},
    {'text': '少喝三杯水→代谢效率下降', 'size': BODY, 'color': GRAY, 'space': 6},
    {'text': '觉得自己代谢不好、容易胖？跟这个有关', 'size': BODY, 'color': GRAY, 'space': 6},
]
add_multiline(slide, Inches(1.0), Inches(1.9), Inches(5), Inches(4.0), lines_water)

add_rounded_rect(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(4.5), ORANGE_BG)
add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.06), ACCENT2)
lines_protein = [
    {'text': '🧱 蛋白质', 'size': LARGE, 'bold': True, 'color': ACCENT2, 'space': 6},
    {'text': '需要：73 g/天  |  实际：45~50 g/天', 'size': BODY, 'color': DARK, 'space': 10},
    {'text': '', 'size': 6, 'color': DARK, 'space': 2},
    {'text': '缺口：-23~28 g/天', 'size': LARGE, 'bold': True, 'color': RED, 'space': 12},
    {'text': '蛋白质是细胞修复的原材料', 'size': BODY, 'color': GRAY, 'space': 6},
    {'text': '缺口这么大，细胞拿什么修自己？', 'size': BODY, 'bold': True, 'color': DARK, 'space': 6},
]
add_multiline(slide, Inches(7.4), Inches(1.9), Inches(5), Inches(4.0), lines_protein)

add_shape(slide, Inches(0.6), Inches(6.4), Inches(12), Inches(0.6), RGBColor(0xFD, 0xE0, 0xE0))
add_textbox(slide, Inches(1.0), Inches(6.44), Inches(11), Inches(0.5),
            '⚠️ 不是偶尔一天——是天天缺、365天都在缺！一天两天扛得住，一年两年呢？十年呢？',
            font_size=BODY, bold=True, color=RED, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# P13 · 基础四件套
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '解决方案——基础四件套')

products = [
    ('🥛 蛋白粉', '细胞修复原材料\n每天几勺补上蛋白质缺口', ACCENT, BLUE_BG),
    ('🍊 维生素C', '抗氧化·胶原合成·免疫力\n人体不能自己合成', ACCENT2, ORANGE_BG),
    ('💊 男士活力组合', '多维+基础鱼油1g+奶蓟草护肝\n一片搞定', GREEN, GREEN_BG),
    ('🐟 高纯鱼油', '4粒/天=1.38g纯Omega-3\n细胞膜的清洁工', RGBColor(0x8E, 0x44, 0xAD), RGBColor(0xF3, 0xE8, 0xFF)),
]
for i, (name, desc, color, bg) in enumerate(products):
    x = Inches(0.4 + i * 3.25)
    add_rounded_rect(slide, x, Inches(1.6), Inches(2.9), Inches(3.5), bg)
    add_shape(slide, x, Inches(1.6), Inches(2.9), Inches(0.06), color)
    add_textbox(slide, x + Inches(0.15), Inches(1.8), Inches(2.6), Inches(0.5),
                name, font_size=SUBTITLE, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    txBox = add_textbox(slide, x + Inches(0.15), Inches(2.5), Inches(2.6), Inches(2.5),
                        '', font_size=SMALL)
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, ll in enumerate(desc.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = ll
        p.font.size = Pt(BODY)
        p.font.color.rgb = DARK
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

# 四件套详情（底部）
add_rounded_rect(slide, Inches(0.4), Inches(5.3), Inches(12.5), Inches(1.8), CARD_BG)
add_shape(slide, Inches(0.4), Inches(5.3), Inches(12.5), Inches(0.05), ACCENT)
details = [
    {'text': '蛋白粉 + 维C → 把基础的"砖瓦"给够', 'size': BODY, 'bold': True, 'color': ACCENT, 'space': 4},
    {'text': '男士活力组合 → 代谢通路打通，奶蓟草护肝  |  鱼油 → 细胞膜洗干净，营养进得去、垃圾出得来', 'size': SMALL, 'color': DARK, 'space': 6},
    {'text': '肝好，身体才能好——这是调理的起点，不是终点', 'size': BODY, 'bold': True, 'color': ACCENT2, 'space': 6},
]
add_multiline(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5), details)

# ══════════════════════════════════════════════════════════════
# P14 · 逻辑链条
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '为什么是这四样？——逻辑链条')

flow = [
    ('四件套', ACCENT),
    ('基础代谢调好', ACCENT2),
    ('肝功能恢复正常', GREEN),
    ('全身细胞修复启动', ACCENT),
]
for i, (text, color) in enumerate(flow):
    x = Inches(0.6 + i * 3.2)
    add_rounded_rect(slide, x, Inches(1.6), Inches(2.6), Inches(1.0), color)
    add_textbox(slide, x, Inches(1.7), Inches(2.6), Inches(0.8),
                text, font_size=BODY, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < len(flow) - 1:
        add_textbox(slide, x + Inches(2.6), Inches(1.7), Inches(0.6), Inches(0.8),
                    '→', font_size=HUGE, bold=True, color=ACCENT2, alignment=PP_ALIGN.CENTER)

roles_p14 = [
    ('蛋白粉 + 维C', '把基础的"砖瓦"给够', ACCENT),
    ('男士活力组合', '代谢通路打通，奶蓟草护肝', ACCENT2),
    ('鱼油', '细胞膜洗干净，营养进得去、垃圾出得来', GREEN),
]
for i, (title, desc, color) in enumerate(roles_p14):
    y = Inches(3.0 + i * 1.3)
    add_rounded_rect(slide, Inches(0.6), y, Inches(12), Inches(1.0), CARD_BG)
    add_shape(slide, Inches(0.6), y, Inches(12), Inches(0.04), color)
    add_textbox(slide, Inches(1.0), y + Inches(0.1), Inches(4), Inches(0.6),
                title, font_size=BODY, bold=True, color=color)
    add_textbox(slide, Inches(5.5), y + Inches(0.1), Inches(7), Inches(0.6),
                desc, font_size=BODY, color=DARK)

add_shape(slide, Inches(0.6), Inches(7.0), Inches(12), Inches(0.4), ACCENT)
add_textbox(slide, Inches(1.0), Inches(7.02), Inches(11), Inches(0.35),
            '肝好，身体才能好——这是调理的起点，不是终点',
            font_size=BODY, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# P15 · 值不值？
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '值不值？——算一算健康这笔账')

add_rounded_rect(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.5), BLUE_BG)
add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.06), ACCENT)
lines_cost = [
    {'text': '基础四件套', 'size': LARGE, 'bold': True, 'color': ACCENT, 'space': 6},
    {'text': '月花费：约 1,400 ~ 2,200 元', 'size': SUBTITLE, 'color': DARK, 'space': 8},
    {'text': '日均：约 50 ~ 70 元', 'size': BODY, 'color': DARK, 'space': 6},
]
add_multiline(slide, Inches(1.0), Inches(1.8), Inches(5), Inches(2.0), lines_cost)

add_rounded_rect(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(2.5), ORANGE_BG)
add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.06), ACCENT2)
lines_cmp = [
    {'text': '你每天在花什么？', 'size': LARGE, 'bold': True, 'color': ACCENT2, 'space': 6},
    {'text': '🚬 一包烟 25元/天 → 750元/月', 'size': BODY, 'color': DARK, 'space': 6},
    {'text': '🧋 一杯奶茶 15元/天 → 450元/月', 'size': BODY, 'color': DARK, 'space': 4},
    {'text': '🍔 一顿外卖加个鸡腿 15元 → 450元/月', 'size': BODY, 'color': DARK, 'space': 4},
]
add_multiline(slide, Inches(7.4), Inches(1.8), Inches(5), Inches(2.0), lines_cmp)

add_rounded_rect(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(2.5), RGBColor(0x0F, 0x17, 0x2A))
lines_value = [
    {'text': '花在手机、衣服、吃喝上的钱，花完就没了', 'size': BODY, 'color': RGBColor(0xAA, 0xBB, 0xDD), 'space': 6},
    {'text': '花在自己细胞上的钱——让你少生病、精神好、老得慢', 'size': SUBTITLE, 'bold': True, 'color': WHITE, 'space': 8},
    {'text': '', 'size': 6, 'color': WHITE, 'space': 2},
    {'text': '这不是消费，是投资', 'size': LARGE, 'bold': True, 'color': ACCENT2, 'space': 6},
]
add_multiline(slide, Inches(1.0), Inches(4.8), Inches(11), Inches(2.0), lines_value)

# ══════════════════════════════════════════════════════════════
# P16 · 长客会
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '怎么花最少的钱？——长客会计划')

# 三列对比
add_rounded_rect(slide, Inches(0.6), Inches(1.6), Inches(3.8), Inches(2.5), BLUE_BG)
add_shape(slide, Inches(0.6), Inches(1.6), Inches(3.8), Inches(0.06), ACCENT)
lines_fish = [
    {'text': '🐟 鱼油算账', 'size': SUBTITLE, 'bold': True, 'color': ACCENT, 'space': 4},
    {'text': '430元/瓶 · 135粒', 'size': SMALL, 'color': DARK, 'space': 6},
    {'text': '4粒/天 → 12.7元/天', 'size': BODY, 'bold': True, 'color': DARK, 'space': 6},
    {'text': '≈ 383元/月', 'size': LARGE, 'bold': True, 'color': ACCENT2, 'space': 6},
]
add_multiline(slide, Inches(0.9), Inches(1.8), Inches(3.2), Inches(2.0), lines_fish)

add_rounded_rect(slide, Inches(4.8), Inches(1.6), Inches(3.8), Inches(2.5), ORANGE_BG)
add_shape(slide, Inches(4.8), Inches(1.6), Inches(3.8), Inches(0.06), ACCENT2)
lines_total = [
    {'text': '📦 四件套正常', 'size': SUBTITLE, 'bold': True, 'color': ACCENT2, 'space': 4},
    {'text': '蛋白粉560+维C347', 'size': SMALL, 'color': DARK, 'space': 4},
    {'text': '+活力组合868+鱼油383', 'size': SMALL, 'color': DARK, 'space': 2},
    {'text': '≈ 2,158元/月', 'size': LARGE, 'bold': True, 'color': RED, 'space': 6},
]
add_multiline(slide, Inches(5.1), Inches(1.8), Inches(3.2), Inches(2.0), lines_total)

add_rounded_rect(slide, Inches(9.0), Inches(1.6), Inches(3.8), Inches(2.5), GREEN_BG)
add_shape(slide, Inches(9.0), Inches(1.6), Inches(3.8), Inches(0.06), GREEN)
lines_changke = [
    {'text': '🎯 长客会', 'size': SUBTITLE, 'bold': True, 'color': GREEN, 'space': 4},
    {'text': '1+N方案', 'size': BODY, 'color': DARK, 'space': 6},
    {'text': '≈ 907元/月', 'size': HUGE, 'bold': True, 'color': GREEN, 'space': 6},
    {'text': '≈ 3.8折', 'size': LARGE, 'bold': True, 'color': ACCENT2, 'space': 4},
]
add_multiline(slide, Inches(9.3), Inches(1.8), Inches(3.2), Inches(2.0), lines_changke)

# 省钱总结
add_rounded_rect(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(2.5), RGBColor(0x0F, 0x17, 0x2A))
lines_save = [
    {'text': '一年能省多少钱？', 'size': LARGE, 'bold': True, 'color': WHITE, 'space': 8},
    {'text': '正常买一年：25,896 元  →  长客会后省下一万三以上', 'size': BODY, 'color': RGBColor(0xAA, 0xBB, 0xDD), 'space': 8},
    {'text': '', 'size': 6, 'color': WHITE, 'space': 2},
    {'text': '省下的钱给自己买个好点的保险，或者给家人也配一套', 'size': BODY, 'color': RGBColor(0xAA, 0xBB, 0xDD), 'space': 6},
    {'text': '分享给朋友，朋友也签了长客会——你还有收入', 'size': SUBTITLE, 'bold': True, 'color': GREEN, 'space': 8},
]
add_multiline(slide, Inches(1.0), Inches(4.7), Inches(11), Inches(2.0), lines_save)

# ══════════════════════════════════════════════════════════════
# P17 · 四句话签约法
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '四句话签约法')

steps_p17 = [
    ('① 戳痛点', '你花几千买保健品觉得贵\n但花几百看病你就愿意了？', ACCENT, BLUE_BG),
    ('② 给方案', '不需要让他多花钱\n反而帮他省钱', ACCENT2, ORANGE_BG),
    ('③ 讲故事', '不说理论，说真实案例\n故事比道理有力量一百倍', GREEN, GREEN_BG),
    ('④ 造紧迫', '给一个现在下单的理由\n先签3个月试试，不合适可以退', RGBColor(0x8E, 0x44, 0xAD), RGBColor(0xF3, 0xE8, 0xFF)),
]
for i, (title, desc, color, bg) in enumerate(steps_p17):
    x = Inches(0.4 + i * 3.25)
    add_rounded_rect(slide, x, Inches(1.6), Inches(2.9), Inches(3.5), bg)
    add_shape(slide, x, Inches(1.6), Inches(2.9), Inches(0.06), color)
    add_textbox(slide, x + Inches(0.15), Inches(1.8), Inches(2.6), Inches(0.5),
                title, font_size=LARGE, bold=True, color=color, alignment=PP_ALIGN.CENTER)
    txBox = add_textbox(slide, x + Inches(0.15), Inches(2.5), Inches(2.6), Inches(2.5),
                        '', font_size=SMALL)
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, ll in enumerate(desc.split('\n')):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = ll
        p.font.size = Pt(BODY)
        p.font.color.rgb = DARK
        p.font.name = 'Microsoft YaHei'
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(8)

# ══════════════════════════════════════════════════════════════
# P18 · 完整话术
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '四句话签约·完整话术')

add_rounded_rect(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(5.5), RGBColor(0x0F, 0x17, 0x2A))
add_shape(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.06), ACCENT2)

script_p18 = [
    {'text': '📝 顺序不能乱：先戳痛点再给方案，先讲故事再造紧迫', 'size': SUBTITLE, 'bold': True, 'color': ACCENT2, 'space': 10},
    {'text': '', 'size': 8, 'color': WHITE, 'space': 2},
    {'text': '❶ 戳痛点："你花几千买保健品觉得贵，但花几百看病你就愿意了？"', 'size': BODY, 'color': RGBColor(0xCC, 0xDD, 0xFF), 'space': 8},
    {'text': '❷ 给方案："有方案让你省心省钱还能赚回更多，愿意了解吗？"', 'size': BODY, 'color': RGBColor(0xCC, 0xDD, 0xFF), 'space': 8},
    {'text': '❸ 讲故事："我有个客户，用这套方法三个月影响了10个人，月入4,000"', 'size': BODY, 'color': RGBColor(0xCC, 0xDD, 0xFF), 'space': 8},
    {'text': '❹ 造紧迫："现在有优惠活动，先签3个月试试，不合适可以退"', 'size': BODY, 'color': RGBColor(0xCC, 0xDD, 0xFF), 'space': 8},
    {'text': '', 'size': 8, 'color': WHITE, 'space': 2},
    {'text': '💡 核心不是"卖东西"，是帮人看到他自己没看到的需求', 'size': SUBTITLE, 'bold': True, 'color': ACCENT2, 'space': 8},
]
add_multiline(slide, Inches(1.0), Inches(1.9), Inches(11), Inches(5.0), script_p18)

# ══════════════════════════════════════════════════════════════
# P19 · 两个选择
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_page_title(slide, '今晚，请你做一个选择')

add_rounded_rect(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.0), CARD_BG)
add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.06), GRAY)
lines_c1 = [
    {'text': '选择一', 'size': LARGE, 'bold': True, 'color': GRAY, 'space': 6},
    {'text': '听完就过去了', 'size': SUBTITLE, 'color': DARK, 'space': 8},
    {'text': '回去该吃吃该喝喝', 'size': BODY, 'color': DARK, 'space': 6},
    {'text': '身体继续每天缺、天天缺', 'size': BODY, 'color': DARK, 'space': 6},
    {'text': '明年体检，箭头更多了', 'size': BODY, 'bold': True, 'color': RED, 'space': 10},
    {'text': '然后你可能会想：', 'size': BODY, 'color': GRAY, 'space': 4},
    {'text': '"去年那个人说的好像有点道理"', 'size': BODY, 'color': GRAY, 'space': 4},
    {'text': '——但已经又过去一年了', 'size': BODY, 'bold': True, 'color': RED, 'space': 6},
]
add_multiline(slide, Inches(1.0), Inches(1.8), Inches(5), Inches(4.5), lines_c1)

add_rounded_rect(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.0), BLUE_BG)
add_shape(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.06), ACCENT)
lines_c2 = [
    {'text': '选择二', 'size': LARGE, 'bold': True, 'color': ACCENT, 'space': 6},
    {'text': '从今天开始改变', 'size': SUBTITLE, 'color': DARK, 'space': 8},
    {'text': '给细胞它需要的东西', 'size': BODY, 'color': DARK, 'space': 6},
    {'text': '不是一下子改变所有习惯', 'size': BODY, 'color': DARK, 'space': 6},
    {'text': '而是从最简单的一步开始', 'size': BODY, 'color': DARK, 'space': 6},
    {'text': '', 'size': 6, 'color': DARK, 'space': 2},
    {'text': '✅ 先把每天缺少的', 'size': SUBTITLE, 'bold': True, 'color': ACCENT, 'space': 6},
    {'text': '✅ 蛋白质和维生素补上', 'size': SUBTITLE, 'bold': True, 'color': ACCENT, 'space': 6},
]
add_multiline(slide, Inches(7.4), Inches(1.8), Inches(5), Inches(4.5), lines_c2)

add_textbox(slide, Inches(0.6), Inches(6.9), Inches(12), Inches(0.5),
            '你的身体是你这辈子最贵的资产——它每天在帮你修复、代谢、免疫，但你需要给它足够的原材料',
            font_size=SMALL, color=GRAY, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# P20 · 行动号召
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.25), H, ACCENT)

add_textbox(slide, Inches(1.2), Inches(1.2), Inches(10), Inches(1.0),
            '想了解自己每天缺什么？', font_size=TITLE, bold=True, color=WHITE)

add_rounded_rect(slide, Inches(1.2), Inches(2.8), Inches(11), Inches(2.0), ACCENT)
add_textbox(slide, Inches(1.8), Inches(3.0), Inches(9.5), Inches(1.6),
            '📩 私信我\n我帮你做个人营养评估，适合什么方案一步到位',
            font_size=LARGE, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1.2), Inches(5.3), Inches(11), Inches(1.5),
            '健康的认知差就是财富差\n\n细胞健康，人就健康',
            font_size=SUBTITLE, color=RGBColor(0xAA, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(1.2), Inches(6.9), Inches(2.5), Inches(0.05), ACCENT2)

# ── 保存 ──
OUT_DIR = '/Users/mac/Documents/知识库/06-输出（成品）/ppt营养修复方案大众版'
os.makedirs(OUT_DIR, exist_ok=True)
out_path = os.path.join(OUT_DIR, '人为什么会生病-大众版.pptx')
prs.save(out_path)
print(f'✅ PPT 已生成：{out_path}')
print(f'   共 {len(prs.slides)} 页')
