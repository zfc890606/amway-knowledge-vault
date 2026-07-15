#!/usr/bin/env python3
"""
Generate Neo-Retro Dev Deck .pptx — 人为什么会生病
复古未来主义·像素信息图·工程笔记本感
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ─── Colors ───
PINK  = RGBColor(0xFF, 0x14, 0x93)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)
CYAN  = RGBColor(0x00, 0xBF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
DARK  = RGBColor(0x33, 0x33, 0x33)
GRAY  = RGBColor(0x88, 0x88, 0x88)
RED   = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x16, 0xA3, 0x4A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color=CREAM):
    """Fill slide background with solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.color.rgb = border_color or BLACK
    shape.line.width = border_width or Pt(3)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=BLACK, alignment=PP_ALIGN.LEFT, font_name='Arial Black'):
    """Add a text box with single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=16, color=DARK, alignment=PP_ALIGN.LEFT, font_name='Arial', line_space=1.4):
    """Add a text box with multiple lines/paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Support tuple (text, size, bold, color)
        if isinstance(line, tuple):
            p.text = line[0]
            p.font.size = Pt(line[1]) if len(line) > 1 else Pt(font_size)
            p.font.bold = line[2] if len(line) > 2 else False
            p.font.color.rgb = line[3] if len(line) > 3 else color
        else:
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_space - 1))
    return txBox

def add_thick_frame(slide, left, top, width, height, shadow_color=None):
    """Thick black bordered rectangle with optional hard shadow."""
    # Main rect
    rect = add_rect(slide, left, top, width, height, fill_color=WHITE, border_color=BLACK, border_width=Pt(4))
    if shadow_color:
        # Shadow rect (behind)
        shadow = add_rect(slide, left + Pt(6), top + Pt(6), width, height, fill_color=shadow_color, border_color=shadow_color, border_width=Pt(0))
        # Move shadow behind
        sp = shadow._element
        sp.getparent().remove(sp)
        rect._element.addprevious(sp)
    return rect

def add_slide_number(slide, num, total=20):
    """Add page number bottom right."""
    add_textbox(slide, Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.4),
                f"{num:02d} / {total:02d}", font_size=10, color=GRAY, alignment=PP_ALIGN.RIGHT, font_name='Courier New')

def add_kicker(slide, text, top=Inches(0.3)):
    """Add monospace kicker at top."""
    add_textbox(slide, Inches(0.5), top, Inches(6), Inches(0.35),
                text, font_size=10, color=GRAY, font_name='Courier New')

def add_h1(slide, text, top=Inches(1.2), color=BLACK, size=44):
    """Huge bold title."""
    add_textbox(slide, Inches(0.7), top, Inches(11), Inches(1.5),
                text, font_size=size, bold=True, color=color, font_name='Arial Black')

def add_h2(slide, text, top=Inches(1.0), color=BLACK, size=32):
    """Bold subtitle."""
    add_textbox(slide, Inches(0.7), top, Inches(11), Inches(1.2),
                text, font_size=size, bold=True, color=color, font_name='Arial Black')

def add_divider_line(slide, top):
    """Thick black horizontal divider."""
    add_rect(slide, Inches(0.5), top, Inches(12.3), Inches(0.04), fill_color=BLACK)

def add_notes(slide, text):
    """Add speaker notes."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text

def new_slide(prs):
    """Create a blank slide."""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


# ═══════════════════════════════════════════════════════════
# P1 · COVER
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
# Large pink accent block top-left
add_rect(sl, Inches(0), Inches(0), Inches(0.6), Inches(7.5), fill_color=PINK)
# Cover title
add_textbox(sl, Inches(1.2), Inches(1.5), Inches(10), Inches(2.5),
            "人为什么\n会生病", font_size=72, bold=True, color=BLACK, font_name='Arial Black')
add_textbox(sl, Inches(1.2), Inches(4.2), Inches(8), Inches(0.6),
            "从隐性饥饿到细胞修复 · 一套让你看懂健康的底层逻辑",
            font_size=18, color=DARK, font_name='Courier New')
# Tags row
tags = ["20 分钟看懂", "底层认知", "解决方案", "长客会签约"]
for i, tag in enumerate(tags):
    x = Inches(1.2) + Inches(2.0) * i
    rect = add_rect(sl, x, Inches(5.2), Inches(1.8), Inches(0.45), fill_color=PINK if i == 0 else WHITE, border_color=BLACK, border_width=Pt(2))
    add_textbox(sl, x, Inches(5.2), Inches(1.8), Inches(0.45), tag,
                font_size=10, bold=True, color=WHITE if i == 0 else BLACK,
                alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_slide_number(sl, 1)
add_notes(sl, "大家好，欢迎来到今天的分享。今天我们不聊复杂的医学理论，只聊一件事——人为什么会生病。我相信在座的各位每年都做体检，但每年体检报告上那些上上下下的箭头，到底意味着什么？医生说没事，注意饮食多运动，但到底怎么注意？没人告诉你。今天我用20分钟，带你从细胞的角度看懂健康的底层逻辑。听完你就知道：身体到底缺了什么、应该怎么补、以及怎么用最少的钱把健康管好。")


# ═══════════════════════════════════════════════════════════
# P2 · 扎心事实
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 01 · 扎心的事实")
add_h2(sl, "一个很多人都有的困惑")
# Pain card
frame = add_thick_frame(sl, Inches(2.5), Inches(2.5), Inches(8.3), Inches(3.5), shadow_color=RED)
add_textbox(sl, Inches(3.0), Inches(2.8), Inches(7.3), Inches(0.6),
            "每年体检，查出箭头", font_size=24, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(3.0), Inches(3.5), Inches(7.3), Inches(0.4),
            "↓", font_size=20, color=RED, alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_textbox(sl, Inches(3.0), Inches(3.9), Inches(7.3), Inches(0.5),
            '医生说："没事，注意饮食多运动"', font_size=20, color=DARK, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(3.0), Inches(4.5), Inches(7.3), Inches(0.4),
            "↓", font_size=20, color=RED, alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_textbox(sl, Inches(3.0), Inches(4.9), Inches(7.3), Inches(0.8),
            "到底怎么注意？\n没人告诉你。", font_size=22, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
add_slide_number(sl, 2)
add_notes(sl, "来，我问一个问题——你每年体检完，看到报告上那些箭头，什么感觉？我问过很多人，十个里面八个说焦虑，但不知道怎么办。拿去问医生，医生翻一翻说没什么大问题，注意饮食多运动——然后呢？到底怎么注意？吃什么？吃多少？没人告诉你。于是你继续该吃吃该喝喝，第二年体检，箭头更多了。这不是你的问题，是整个体系里缺了一环——没人从底层告诉你身体到底需要什么。今天这一课，就是补上这一环的。")


# ═══════════════════════════════════════════════════════════
# P3 · 今晚目标
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 01 · 今晚的目标")
add_h2(sl, "听完你就知道")
goals = [
    ("①", "身体到底缺了什么", "为什么一日三餐吃得饱饱的，\n身体还是会出问题"),
    ("②", "应该怎么补", "三件套方案，\n把每天缺的补回来"),
    ("③", "怎么省钱吃、甚至赚钱", "了解长客会，\n花更少的钱买更好的营养"),
]
for i, (num, title, desc) in enumerate(goals):
    x = Inches(0.8) + Inches(4.1) * i
    add_thick_frame(sl, x, Inches(2.5), Inches(3.7), Inches(3.8), shadow_color=PINK if i == 0 else None)
    add_textbox(sl, x, Inches(2.8), Inches(3.7), Inches(0.7),
                num, font_size=36, bold=True, color=PINK, alignment=PP_ALIGN.CENTER, font_name='Courier New')
    add_textbox(sl, x, Inches(3.5), Inches(3.7), Inches(0.6),
                title, font_size=18, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_textbox(sl, x, Inches(4.3), Inches(3.7), Inches(1.5),
                desc, font_size=13, color=DARK, alignment=PP_ALIGN.CENTER)
add_slide_number(sl, 3)
add_notes(sl, "好，我先给大家一个今晚的地图。第一，你会知道身体到底缺了什么。第二，你会知道应该怎么补。第三，也是大家最关心的——怎么省钱吃，甚至还能赚钱。安利的长客会计划，如果你还不知道，今晚听完可能能帮你一年省好几千。好，目标清楚了，我们开始进入正题。")


# ═══════════════════════════════════════════════════════════
# P4 · 从细胞说起
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 02 · 底层认知")
add_h2(sl, '一切从细胞说起', top=Inches(0.8), color=PINK)
# Hierarchy blocks
levels = [
    ("细胞", PINK), ("组织", RGBColor(0xD9, 0x46, 0xEF)),
    ("器官", CYAN), ("系统", YELLOW), ("人", BLACK),
]
for i, (name, clr) in enumerate(levels):
    x = Inches(0.7) + Inches(2.4) * i
    rect = add_rect(sl, x, Inches(2.5), Inches(1.8), Inches(0.8), fill_color=clr, border_color=BLACK, border_width=Pt(3))
    text_color = WHITE if clr != YELLOW else BLACK
    add_textbox(sl, x, Inches(2.6), Inches(1.8), Inches(0.6),
                name, font_size=18, bold=True, color=text_color, alignment=PP_ALIGN.CENTER, font_name='Courier New')
    if i < len(levels) - 1:
        add_textbox(sl, x + Inches(1.85), Inches(2.55), Inches(0.5), Inches(0.5),
                    "→", font_size=20, bold=True, font_name='Courier New')
# Formula box
fx = Inches(3.5)
frame = add_rect(sl, fx, Inches(4.5), Inches(6.3), Inches(1.3), fill_color=WHITE, border_color=BLACK, border_width=Pt(4))
add_textbox(sl, fx, Inches(4.55), Inches(6.3), Inches(0.3),
            "核心公式", font_size=10, color=PINK, alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_textbox(sl, fx, Inches(4.9), Inches(6.3), Inches(0.7),
            "细胞健康  =  人健康", font_size=22, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
add_slide_number(sl, 4)
add_notes(sl, "好，我们从头说起。人体是由细胞构成的。细胞组成组织，组织组成器官，器官组成系统，系统组成人。关键是这个公式——细胞健康等于人健康。这不是一句口号，是现代医学的底层共识。如果你每个细胞都是健康的，那你的组织就是健康的，器官就是健康的，人就是健康的。反过来，如果细胞出了问题，你吃再多药都只是治标不治本。")


# ═══════════════════════════════════════════════════════════
# P5 · 7大营养素
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 02 · 底层认知")
add_h2(sl, '细胞的"原材料"——7 大营养素')
nutrients = ["水", "蛋白质", "维生素", "矿物质", "碳水", "脂肪", "膳食纤维"]
for i, n in enumerate(nutrients):
    x = Inches(0.5) + Inches(1.8) * i
    add_rect(sl, x, Inches(2.8), Inches(1.5), Inches(1.5), fill_color=WHITE, border_color=BLACK, border_width=Pt(3))
    add_textbox(sl, x, Inches(3.5), Inches(1.5), Inches(0.5),
                n, font_size=16, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_textbox(sl, Inches(1), Inches(5.0), Inches(10), Inches(0.5),
            "这些就是细胞正常工作的原材料", font_size=18, color=DARK, alignment=PP_ALIGN.CENTER)
add_slide_number(sl, 5)
add_notes(sl, "细胞要正常工作，需要7大类营养素：水、蛋白质、维生素、矿物质、碳水、脂肪、膳食纤维。你可能会想这些我知道啊，但关键是你每天吃进去的量，够不够？水是身体的溶剂，蛋白质是修房子的砖，维生素和矿物质是工具和钥匙。这7样东西缺任何一样，细胞就没法正常工作。")


# ═══════════════════════════════════════════════════════════
# P6 · 一句话逻辑
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 02 · 底层认知")
add_textbox(sl, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
            '原材料不够，\n细胞就没法正常干活。', font_size=40, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER, font_name='Arial Black')
add_textbox(sl, Inches(1.5), Inches(3.8), Inches(10), Inches(1.2),
            '人就会出问题。', font_size=48, bold=True, color=PINK, alignment=PP_ALIGN.CENTER, font_name='Arial Black')
add_rect(sl, Inches(5.5), Inches(5.2), Inches(2.3), Inches(0.06), fill_color=BLACK)
add_textbox(sl, Inches(1.5), Inches(5.4), Inches(10), Inches(0.5),
            '—— 就这么简单', font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_slide_number(sl, 6)
add_notes(sl, "原材料不够，细胞就没法正常干活。人就会出问题。我希望大家把这句话记住。整个健康管理的底层逻辑，就这一句话。你的身体每分每秒都在用营养素做修复、做代谢、做免疫。原材料够了，身体自己会修自己。原材料不够，修不好，问题就一点点积累。那问题来了——我们每天的砖够吗？")


# ═══════════════════════════════════════════════════════════
# P7 · 吃≠吃饱
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 03 · 核心概念：隐性饥饿")
add_h2(sl, "你吃进去的，跟细胞需要的\n是两码事")
# Left card
add_thick_frame(sl, Inches(0.8), Inches(3.5), Inches(4.5), Inches(3.0), shadow_color=RED)
add_textbox(sl, Inches(0.8), Inches(3.6), Inches(4.5), Inches(0.3),
            "你以为", font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_textbox(sl, Inches(0.8), Inches(4.3), Inches(4.5), Inches(1.0),
            "一日三餐\n吃得饱饱的", font_size=16, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
# VS
add_rect(sl, Inches(5.8), Inches(4.5), Inches(1.0), Inches(1.0), fill_color=BLACK)
add_textbox(sl, Inches(5.8), Inches(4.6), Inches(1.0), Inches(0.8),
            "VS", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, font_name='Courier New')
# Right card
add_thick_frame(sl, Inches(7.3), Inches(3.5), Inches(5.2), Inches(3.0), shadow_color=GREEN)
add_textbox(sl, Inches(7.3), Inches(3.6), Inches(5.2), Inches(0.3),
            "实际上", font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_textbox(sl, Inches(7.3), Inches(4.3), Inches(5.2), Inches(1.5),
            "碳水够够的\n蛋白质/维生素/矿物质\n远远不够", font_size=16, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
add_slide_number(sl, 7)
add_notes(sl, "很多人跟我说我一天三顿饭，吃得很好啊。但你吃进去的，跟细胞需要的，是两码事。左边是你以为的吃饱了——米饭面条包子饺子，全是碳水。右边是细胞真正需要的——优质蛋白、维生素、矿物质。中国人的饮食习惯，一碗米饭配两个菜，一顿可能就20-25克蛋白质。一天三顿才60-75克。所以吃饱和吃够是两回事。")


# ═══════════════════════════════════════════════════════════
# P8 · 隐性饥饿
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 03 · 核心概念：隐性饥饿")
add_h2(sl, "隐性饥饿")
add_textbox(sl, Inches(1.5), Inches(1.8), Inches(10), Inches(0.4),
            "胃饱了，细胞还饿着", font_size=18, color=CYAN, font_name='Courier New')
add_textbox(sl, Inches(0.8), Inches(3.0), Inches(11), Inches(0.7),
            "每天缺一点，感觉不到。", font_size=28, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(0.8), Inches(3.8), Inches(11), Inches(0.4),
            "↓", font_size=22, color=PINK, alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_textbox(sl, Inches(0.8), Inches(4.3), Inches(11), Inches(0.8),
            '直到体检报告出来——\n问题来了', font_size=24, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
# Analogy box
add_thick_frame(sl, Inches(2.5), Inches(5.5), Inches(8.3), Inches(1.3), shadow_color=CYAN)
add_textbox(sl, Inches(3.0), Inches(5.6), Inches(7.3), Inches(0.3),
            "就像水蛭吸血：每天吸走一点点，你完全没感觉。直到某天你站起来，发现头晕眼花——已经贫血了。",
            font_size=14, color=DARK, alignment=PP_ALIGN.LEFT, font_name='Arial')
add_slide_number(sl, 8)
add_notes(sl, "隐性饥饿——这是今天最重要的一个概念。世界卫生组织定义：隐性饥饿是指微量营养素摄入不足，但你可能完全感觉不到。为什么叫隐性？因为它不疼不痒。它是每天缺一点点——水少喝200毫升，蛋白质少吃10克——一天两天没事，一年两年，问题就出来了。就像水蛭吸血，每天吸一点点，等你发现的时候——已经出问题了。")


# ═══════════════════════════════════════════════════════════
# P9 · 为什么感觉不到
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 03 · 核心概念：隐性饥饿")
add_textbox(sl, Inches(1.0), Inches(2.0), Inches(11), Inches(1.5),
            "慢性病\n不是突然发生的", font_size=52, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER, font_name='Arial Black')
add_textbox(sl, Inches(1.0), Inches(3.8), Inches(11), Inches(0.4),
            "— — —", font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_textbox(sl, Inches(1.0), Inches(4.3), Inches(11), Inches(1.5),
            '是突然发现的', font_size=56, bold=True, color=PINK, alignment=PP_ALIGN.CENTER, font_name='Arial Black')
# Underline box
add_rect(sl, Inches(4.5), Inches(6.0), Inches(4.3), Inches(0.08), fill_color=YELLOW)
add_slide_number(sl, 9)
add_notes(sl, "这句话我希望大家刻在脑子里——慢性病不是突然发生的，是突然发现的。高血压不是昨天才高的，是血管壁慢慢变硬了十年。糖尿病不是上个月才得的，是胰岛细胞累趴下了好几年。脂肪肝不是这周才有的，是肝细胞被脂肪包围了好几年。你感觉不到，因为身体的代偿能力非常强。肝坏掉70%还能正常工作。但代偿不是治疗——它只是把问题往后推。")


# ═══════════════════════════════════════════════════════════
# P10 · 生命器官理论
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 04 · 生命器官理论")
add_h2(sl, "营养不够时，身体会取舍")
# Vital organs
add_thick_frame(sl, Inches(0.8), Inches(2.5), Inches(5.2), Inches(4.0), shadow_color=GREEN)
add_textbox(sl, Inches(0.8), Inches(2.6), Inches(5.2), Inches(0.7),
            "生命器官\n优先供给", font_size=18, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
vitals = ["心脏", "肺", "肝脏", "肾脏", "大脑"]
for i, v in enumerate(vitals):
    add_rect(sl, Inches(1.2), Inches(3.5 + i * 0.55), Inches(4.4), Inches(0.45), fill_color=CREAM, border_color=BLACK, border_width=Pt(2))
    add_textbox(sl, Inches(1.2), Inches(3.5 + i * 0.55), Inches(4.4), Inches(0.45),
                f"  {v}", font_size=14, bold=True, color=BLACK, font_name='Courier New', alignment=PP_ALIGN.CENTER)
# Divider
add_rect(sl, Inches(6.4), Inches(2.5), Inches(0.06), Inches(4.0), fill_color=BLACK)
# Sacrifice organs
add_thick_frame(sl, Inches(7.0), Inches(2.5), Inches(5.5), Inches(4.0), shadow_color=RED)
add_textbox(sl, Inches(7.0), Inches(2.6), Inches(5.5), Inches(0.7),
            "非生命器官\n被牺牲", font_size=18, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
sacs = ["胃", "乳腺", "甲状腺", "子宫 / 卵巢", "皮肤 / 四肢"]
for i, s in enumerate(sacs):
    c = CREAM if i < 3 else RGBColor(0xFF, 0xF0, 0xF0)
    add_rect(sl, Inches(7.4), Inches(3.5 + i * 0.55), Inches(4.7), Inches(0.45), fill_color=c, border_color=RED, border_width=Pt(2))
    add_textbox(sl, Inches(7.4), Inches(3.5 + i * 0.55), Inches(4.7), Inches(0.45),
                f"  {s}", font_size=14, color=DARK, font_name='Courier New', alignment=PP_ALIGN.CENTER)
add_slide_number(sl, 10)
add_notes(sl, "这里我要讲一个非常重要的理论——生命器官和非生命器官之分。当身体营养不够的时候，它会优先保证心、肝、肺、肾、脑这五个生命器官的供给。那牺牲的是谁？胃、乳腺、甲状腺、子宫、卵巢、皮肤、四肢。所以为什么现在乳腺增生、甲状腺结节这么普遍？不是突然变多了，是它们的营养被挪用了。身体在用你的乳腺和甲状腺的营养，去保你的心和脑。")


# ═══════════════════════════════════════════════════════════
# P11 · 你身上的信号
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 04 · 生命器官理论")
add_h2(sl, "你身上的信号")
add_textbox(sl, Inches(0.7), Inches(1.5), Inches(10), Inches(0.4),
            "不是有病，是原材料不够了", font_size=16, color=CYAN, font_name='Courier New')
signals = ["反复口腔溃疡", "伤口愈合慢", "皮肤松弛、头发枯黄", "乳腺增生",
           "甲状腺结节", "关节不舒服", "总是觉得累", "免疫力差，容易感冒"]
for i, sig in enumerate(signals):
    col = i % 4
    row = i // 4
    x = Inches(0.5) + Inches(3.1) * col
    y = Inches(2.5) + Inches(0.9) * row
    add_rect(sl, x, y, Inches(2.8), Inches(0.65), fill_color=WHITE, border_color=BLACK, border_width=Pt(2))
    add_textbox(sl, x, y, Inches(2.8), Inches(0.65),
                sig, font_size=13, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
# Conclusion
add_rect(sl, Inches(2.0), Inches(5.0), Inches(9.3), Inches(0.65), fill_color=YELLOW, border_color=BLACK, border_width=Pt(2))
add_textbox(sl, Inches(2.0), Inches(5.0), Inches(9.3), Inches(0.65),
            "这些都不是病，是身体在说：原材料不够了！", font_size=18, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
add_slide_number(sl, 11)
add_notes(sl, "说到这里，大家可能会想那我怎么知道自己缺不缺？好，我列几个信号：口腔溃疡反复发作、伤口愈合慢、皮肤松弛头发枯黄、乳腺增生、甲状腺结节、关节不舒服、总是觉得累、免疫力差容易感冒——这里面你中了几个？很多人以为这是亚健康，不对。这些信号的本质，是你的非生命器官在喊救命。身体把你的原材料都调去保心肝肺肾脑了，你的皮肤、头发、乳腺、甲状腺就分不到足够营养了。")


# ═══════════════════════════════════════════════════════════
# P12 · 算账
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 04 · 生命器官理论")
add_h2(sl, "算笔账：你每天缺多少")
# Water card
for idx, (icon, title, need_val, need_lbl, actual_val, actual_lbl, deficit, deficit_lbl, shadow_clr) in enumerate([
    ("水",  "水",  "1,800", "ml/天", "1,000-1,200", "ml/天", "-600~800", "ml/天", CYAN),
    ("蛋白质", "蛋白质", "73", "g/天", "45-50", "g/天", "-23~28", "g/天", PINK),
]):
    x = Inches(0.8) + Inches(6.2) * idx
    add_thick_frame(sl, x, Inches(2.3), Inches(5.6), Inches(3.8), shadow_color=shadow_clr)
    add_textbox(sl, x, Inches(2.4), Inches(5.6), Inches(0.5), title,
                font_size=18, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
    items = [
        ("需要", need_val, need_lbl, BLACK),
        ("实际", actual_val, actual_lbl, RGBColor(0xD9, 0x77, 0x06)),
        ("缺口", deficit, deficit_lbl, RED),
    ]
    for j, (lbl, val, unit, clr) in enumerate(items):
        y = Inches(3.2 + j * 0.9)
        add_textbox(sl, x + Inches(0.3), y, Inches(2.5), Inches(0.4),
                    lbl, font_size=14, color=DARK, font_name='Courier New')
        add_textbox(sl, x + Inches(2.8), y, Inches(2.5), Inches(0.4),
                    f"{val} {unit}", font_size=16, bold=True, color=clr, alignment=PP_ALIGN.RIGHT, font_name='Courier New')
        if j < 2:
            add_rect(sl, x + Inches(0.3), y + Inches(0.45), Inches(5.0), Inches(0.01), fill_color=RGBColor(0xDD, 0xDD, 0xDD))
        if j == 2 and idx == 0:
            add_rect(sl, x + Inches(0.3), y - Inches(0.05), Inches(5.0), Inches(0.02), fill_color=RED, border_color=RED)
add_textbox(sl, Inches(1), Inches(6.4), Inches(11), Inches(0.5),
            "每天缺，天天缺，365 天都在缺", font_size=18, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
add_slide_number(sl, 12)
add_notes(sl, "好，我们来算一笔具体的账。先说水。60公斤的人每天需要1800毫升水，但大部分人每天只喝1000到1200毫升——差了将近600到800毫升。再说蛋白质。60公斤的人每天需要约73克蛋白质，但实际摄入只有45到50克——缺口23到28克，将近三分之一。而且关键的是——这不是偶尔一天，是天天缺！")


# ═══════════════════════════════════════════════════════════
# P13 · 基础三件套
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 05 · 解决方案")
add_h2(sl, "解决方案：基础三件套")
products = [
    ("蛋白粉", "细胞修复原材料\n补充每天缺少的20+克蛋白质", PINK),
    ("维生素C", "抗氧化 · 合成胶原蛋白\n提高免疫力", YELLOW),
    ("活力组合", "多维 + 鱼油 + 护肝片\n全面补充微量营养", CYAN),
]
for i, (name, desc, clr) in enumerate(products):
    x = Inches(0.8) + Inches(4.1) * i
    add_thick_frame(sl, x, Inches(2.5), Inches(3.7), Inches(3.5), shadow_color=clr)
    add_rect(sl, x + Inches(1.2), Inches(2.8), Inches(1.3), Inches(1.3), fill_color=clr, border_color=BLACK, border_width=Pt(2))
    add_textbox(sl, x, Inches(4.3), Inches(3.7), Inches(0.5),
                name, font_size=18, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_textbox(sl, x, Inches(4.9), Inches(3.7), Inches(0.8),
                desc, font_size=13, color=DARK, alignment=PP_ALIGN.CENTER)
add_slide_number(sl, 13)
add_notes(sl, "好，现在来讲解决方案。我推荐的基础三件套——不是让你买一大堆瓶瓶罐罐，而是从最基础的三个缺口开始补。第一，蛋白粉——细胞修复的原材料。第二，维生素C——抗氧化、合成胶原蛋白、提高免疫力。第三，活力组合——多维、鱼油、护肝片，全面补充微量营养。")


# ═══════════════════════════════════════════════════════════
# P14 · 为什么是这三样
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 05 · 解决方案")
add_h2(sl, "为什么是这三样？")
# Flow diagram
flow_items = [
    ("三件套", PINK),
    ("基础代谢\n调好了", YELLOW),
    ("肝功能\n恢复正常", CYAN),
    ("全身细胞\n修复启动", BLACK),
]
for i, (name, clr) in enumerate(flow_items):
    x = Inches(0.8) + Inches(3.1) * i
    rect = add_rect(sl, x, Inches(2.8), Inches(2.5), Inches(1.2), fill_color=clr, border_color=BLACK, border_width=Pt(3))
    txt_color = WHITE if clr != YELLOW else BLACK
    add_textbox(sl, x, Inches(3.1), Inches(2.5), Inches(0.8),
                name, font_size=15, bold=True, color=txt_color, alignment=PP_ALIGN.CENTER, font_name='Courier New')
    if i < len(flow_items) - 1:
        add_textbox(sl, x + Inches(2.55), Inches(3.0), Inches(0.5), Inches(0.5),
                    "→", font_size=22, bold=True, font_name='Courier New')
# Key message
add_thick_frame(sl, Inches(2.5), Inches(4.8), Inches(8.3), Inches(1.3))
add_textbox(sl, Inches(2.8), Inches(4.9), Inches(7.7), Inches(0.5),
            "肝好，身体才能好", font_size=22, bold=True, color=PINK, alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_textbox(sl, Inches(2.8), Inches(5.4), Inches(7.7), Inches(0.5),
            "先把基础代谢和肝功能调好，身体才有能力自己修自己。", font_size=14, color=DARK, alignment=PP_ALIGN.CENTER)
add_slide_number(sl, 14)
add_notes(sl, "你可能会问为什么是这三样？很多人补充营养的顺序错了——一上来就吃钙铁锌硒，但你的基础代谢和肝功能没有调好，吃进去的东西根本吸收不了。三件套的逻辑是：先用蛋白粉和维生素C把基础的砖给够，用活力组合把代谢通路打通。肝是全身的化工厂，几乎所有营养物质的代谢都要经过肝。肝好，身体才能好。")


# ═══════════════════════════════════════════════════════════
# P15 · 值不值
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 05 · 解决方案")
add_textbox(sl, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "值不值？", font_size=52, bold=True, color=PINK, alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_thick_frame(sl, Inches(3.0), Inches(2.8), Inches(7.3), Inches(1.0))
add_textbox(sl, Inches(3.0), Inches(2.85), Inches(7.3), Inches(0.8),
            "900 ~ 1,700  元/月", font_size=36, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_textbox(sl, Inches(1.5), Inches(4.0), Inches(10), Inches(0.8),
            "你觉得自己每天的身体维护，值不值这个钱？", font_size=20, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
# Comparison
add_rect(sl, Inches(3.5), Inches(5.0), Inches(6.3), Inches(1.5), fill_color=WHITE, border_color=BLACK, border_width=Pt(2))
comps = [("一包烟", "25元/天 · 750元/月"), ("一杯奶茶", "15元/天 · 450元/月"), ("三件套", "≈ 30~55元/天")]
for i, (item, price) in enumerate(comps):
    y = Inches(5.1 + i * 0.45)
    add_textbox(sl, Inches(3.8), y, Inches(3.0), Inches(0.4),
                item, font_size=14, color=DARK, font_name='Courier New')
    add_textbox(sl, Inches(6.8), y, Inches(2.8), Inches(0.4),
                price, font_size=14, bold=(i == 2), color=BLACK if i == 2 else DARK,
                alignment=PP_ALIGN.RIGHT, font_name='Courier New')
add_slide_number(sl, 15)
add_notes(sl, "好，大家最关心的问题来了——这一套要多少钱？基础三件套一个月大概900到1700元，一天30到55块钱。对比一下——一包烟25块钱一天，一个月750。一杯奶茶15块钱一天，一个月450。这些钱你花得眼睛都不眨一下。但说到花30到55块钱一天来维护自己的身体，你就犹豫了？这不是消费，是投资。")


# ═══════════════════════════════════════════════════════════
# P16 · 算经济账
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 06 · 长客会签约")
add_h2(sl, "算一笔经济账")
# Price table frame
add_thick_frame(sl, Inches(1.5), Inches(2.3), Inches(10.3), Inches(4.0))
headers = ["方案", "月花费", "折扣"]
rows = [
    ("正常购买", "1,775 元", "—", WHITE, False),
    ("1+N 长客会", "907 元", "≈ 6.5 折 🔥", RGBColor(0xF0, 0xFF, 0xF0), True),
    ("买 6 送 1", "868 元", "更低", RGBColor(0xE8, 0xFF, 0xE8), True),
    ("积分兑换后", "≈ 3.8 折", "最划算 🏆", RGBColor(0xFF, 0xF8, 0xE0), True),
]
# Header
add_rect(sl, Inches(1.5), Inches(2.3), Inches(10.3), Inches(0.5), fill_color=BLACK)
for j, h in enumerate(headers):
    xs = [Inches(1.8), Inches(6.5), Inches(9.0)]
    add_textbox(sl, xs[j], Inches(2.33), Inches(3.0), Inches(0.45),
                h, font_size=11, bold=True, color=WHITE, font_name='Courier New', alignment=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT)
for i, (plan, cost, discount, bg, high) in enumerate(rows):
    y = Inches(2.9 + i * 0.7)
    add_rect(sl, Inches(1.5), y, Inches(10.3), Inches(0.65), fill_color=bg, border_color=BLACK, border_width=Pt(0))
    add_textbox(sl, Inches(1.8), y + Inches(0.05), Inches(4.5), Inches(0.5),
                plan, font_size=15, bold=True if high else False, color=BLACK if high else DARK, font_name='Courier New')
    cost_color = GREEN if high else BLACK
    add_textbox(sl, Inches(6.5), y + Inches(0.05), Inches(2.5), Inches(0.5),
                cost, font_size=16, bold=True, color=cost_color, alignment=PP_ALIGN.RIGHT, font_name='Courier New')
    add_textbox(sl, Inches(9.0), y + Inches(0.05), Inches(2.5), Inches(0.5),
                discount, font_size=11, bold=True, color=BLACK, alignment=PP_ALIGN.RIGHT, font_name='Courier New')
add_slide_number(sl, 16)
add_notes(sl, "好，现在进入大家最期待的部分——怎么花最少的钱把营养补到位。正常买基础三件套一个月大概1775元。但长客会1+N方案只要907元一个月，直接打了6.5折。如果你买6送1，摊下来一个月868元。加上积分兑换——实际到手相当于3.8折。一年下来省一万多！")


# ═══════════════════════════════════════════════════════════
# P17 · 四句话签约法
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 06 · 长客会签约")
add_h2(sl, "四句话签约法")
steps = [
    ("1", "戳痛点", '"花几千买保健品你觉得贵，\n但花几百看病你就愿意了？"'),
    ("2", "给方案", '"现在有个方案，少花钱维护健康，\n把原来该花的钱赚回来。"'),
    ("3", "讲故事", '"我有个客户，用这套方法3个月\n影响了10个人，月入4,000。"'),
    ("4", "造紧迫", '"本月积分翻倍活动，\n先签3个月试试，不合适可以退。"'),
]
for i, (num, label, quote) in enumerate(steps):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6.2) * col
    y = Inches(2.3) + Inches(2.5) * row
    add_thick_frame(sl, x, y, Inches(5.6), Inches(2.0))
    add_rect(sl, x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.6), fill_color=BLACK)
    add_textbox(sl, x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.6),
                num, font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, font_name='Courier New')
    add_textbox(sl, x + Inches(1.0), y + Inches(0.15), Inches(4.3), Inches(0.3),
                label, font_size=12, bold=True, color=PINK, font_name='Courier New')
    add_textbox(sl, x + Inches(1.0), y + Inches(0.6), Inches(4.3), Inches(1.0),
                quote, font_size=14, color=DARK)
add_slide_number(sl, 17)
add_notes(sl, "好，最后分享一个实战技巧——四句话签约法。第一句戳痛点——很多人舍不得花钱在健康上，但愿意花钱看病。第二句给方案——不需要让他多花钱，反而帮他省钱。第三句讲故事——不要说理论，说真实案例。第四句造紧迫——给一个现在下单的理由，同时降低门槛。")


# ═══════════════════════════════════════════════════════════
# P18 · 签约话术
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 06 · 长客会签约")
add_h2(sl, "签约话术全文")
script_lines = [
    ('①', '你想想，你一年花在健康上的钱有多少？花几千买保健品你觉得贵，但花几百看病你就愿意了？这不是省钱，是等着一笔更大的开销。'),
    ('②', '现在有个方案——长客会。签个约，每个月自动配货，比单买便宜35%以上。你少花了钱维护了健康，还能把原来该花的钱赚回来。'),
    ('③', '我有个客户，自己也用这套方案，三个月影响了10个人了解同样的理念，一个月多赚4,000块钱。他不是销售，就是分享自己的体验。'),
    ('④', '这个月入会积分翻倍，你算算你能省多少。先签3个月试试，觉得不合适随时退，没损失。'),
]
add_thick_frame(sl, Inches(1.5), Inches(2.3), Inches(10.3), Inches(4.0))
for i, (num, text) in enumerate(script_lines):
    y = Inches(2.5 + i * 0.9)
    add_textbox(sl, Inches(1.8), y + Inches(0.05), Inches(0.5), Inches(0.4),
                num, font_size=16, bold=True, color=CYAN, font_name='Courier New')
    add_textbox(sl, Inches(2.5), y, Inches(9.0), Inches(0.8),
                text, font_size=13, color=DARK)
    if i < len(script_lines) - 1:
        add_rect(sl, Inches(1.8), y + Inches(0.75), Inches(9.7), Inches(0.02), fill_color=BLACK)
add_slide_number(sl, 18)
add_notes(sl, "这一页是完整的四句话话术，大家可以拍照或者记下来。注意顺序不能乱：先戳痛点再给方案，先讲故事再造紧迫。这个话术的核心不是卖东西，是帮人看到他自己没看到的需求。从细胞需要什么，到你缺什么，到怎么补，到怎么省钱补——这条逻辑链你掌握了，跟任何人聊健康你都有底气。")


# ═══════════════════════════════════════════════════════════
# P19 · 两个选择
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 07 · 结尾")
add_h2(sl, "你现在有两个选择")
# Bad choice
add_thick_frame(sl, Inches(0.8), Inches(2.8), Inches(5.0), Inches(3.2), shadow_color=RED)
add_textbox(sl, Inches(0.8), Inches(2.9), Inches(5.0), Inches(0.7),
            "选择一", font_size=20, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(0.8), Inches(3.8), Inches(5.0), Inches(1.5),
            "听完就过去\n继续每天缺、天天缺\n明年体检箭头更多",
            font_size=16, color=DARK, alignment=PP_ALIGN.CENTER)
# VS
add_rect(sl, Inches(6.0), Inches(3.9), Inches(1.0), Inches(1.0), fill_color=BLACK)
add_textbox(sl, Inches(6.0), Inches(4.0), Inches(1.0), Inches(0.8),
            "OR", font_size=16, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER, font_name='Courier New')
# Good choice
add_thick_frame(sl, Inches(7.3), Inches(2.8), Inches(5.2), Inches(3.2), shadow_color=GREEN)
add_textbox(sl, Inches(7.3), Inches(2.9), Inches(5.2), Inches(0.7),
            "选择二", font_size=20, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(7.3), Inches(3.8), Inches(5.2), Inches(1.5),
            "从现在开始\n给细胞它需要的东西\n让身体自己修自己",
            font_size=16, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
add_slide_number(sl, 19)
add_notes(sl, "好，今天的分享快结束了。临收尾的时候，我想请你做一个选择。选择一——听完就过去了，继续每天缺、天天缺。选择二——从今天开始，给细胞它需要的东西。你的身体是你这辈子最贵的资产，它每天在帮你修复、代谢、免疫——但你需要给它足够的原材料。")


# ═══════════════════════════════════════════════════════════
# P20 · 行动号召
# ═══════════════════════════════════════════════════════════
sl = new_slide(prs)
add_bg(sl, CREAM)
add_kicker(sl, "// PART 07 · 行动号召")
add_textbox(sl, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
            "想了解更多？", font_size=36, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
# CTA card
add_thick_frame(sl, Inches(3.0), Inches(2.8), Inches(7.3), Inches(3.0), shadow_color=PINK)
add_textbox(sl, Inches(3.0), Inches(3.0), Inches(7.3), Inches(0.7),
            "私信我，", font_size=32, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(3.0), Inches(3.8), Inches(7.3), Inches(0.9),
            "我帮你算算", font_size=44, bold=True, color=PINK, alignment=PP_ALIGN.CENTER)
add_textbox(sl, Inches(3.0), Inches(4.8), Inches(7.3), Inches(0.5),
            "你每天缺多少 · 适合什么方案", font_size=16, color=DARK, alignment=PP_ALIGN.CENTER, font_name='Courier New')
# Tags
tags = ["免费评估", "一对一咨询", "3 天试用"]
for i, tag in enumerate(tags):
    x = Inches(3.5 + i * 2.5)
    add_rect(sl, x, Inches(6.0), Inches(2.0), Inches(0.45), fill_color=PINK if i == 0 else WHITE, border_color=BLACK, border_width=Pt(2))
    add_textbox(sl, x, Inches(6.0), Inches(2.0), Inches(0.45),
                tag, font_size=10, bold=True, color=WHITE if i == 0 else BLACK,
                alignment=PP_ALIGN.CENTER, font_name='Courier New')
add_slide_number(sl, 20)
add_notes(sl, "好，今天就分享到这里。如果你觉得今天的内容对你有帮助，想了解自己每天缺什么、适合什么方案——私信我，我帮你做一个个人的营养评估。或者你身边有朋友也有类似的困惑，把这份分享转发给他。记住：细胞健康，人就健康。从今天开始，给你的细胞它需要的东西。")


# ─── Save ───
output_path = "/Users/mac/Documents/知识库/06-输出（成品）/ppt营养修复方案大众版/人为什么会生病-复古未来风.pptx"
prs.save(output_path)
print(f"✅ PPTX saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
