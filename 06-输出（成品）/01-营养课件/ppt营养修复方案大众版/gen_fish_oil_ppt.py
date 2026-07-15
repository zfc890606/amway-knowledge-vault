#!/usr/bin/env python3
"""生成「更新的鱼油」独立PPT——字体大、无动画、纯干货"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 颜色 ──
DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # 深蓝黑
ACCENT  = RGBColor(0x00, 0x6D, 0xAA)   # 科技蓝
ACCENT2 = RGBColor(0xF0, 0x8C, 0x00)   # 橙黄
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF5, 0xF7, 0xFA)
GRAY    = RGBColor(0x66, 0x66, 0x66)
GREEN   = RGBColor(0x2E, 0xCC, 0x71)
RED     = RGBColor(0xE7, 0x4C, 0x3C)
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(2)
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=24, bold=False,
                color=DARK, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_paragraph(tf, text, font_size=22, bold=False, color=DARK,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(8), font_name='Microsoft YaHei'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    return p

# ══════════════════════════════════════════════════════════════
# P1 · 封面
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_DARK)

# 装饰条
add_shape(slide, Inches(0), Inches(0), Inches(0.3), H, ACCENT)

# 主标题
add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.8),
            '更新的鱼油', font_size=72, bold=True, color=WHITE)

# 副标题
add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(1.2),
            '高纯度 Omega-3 · 细胞换膜方案', font_size=36, color=RGBColor(0x8E, 0xC6, 0xFF))

# 底部信息
add_textbox(slide, Inches(1.5), Inches(5.8), Inches(8), Inches(1),
            '纽崔莱™ 新款高纯鱼油 · 345mg/粒 · 细胞级修复', font_size=22, color=GRAY)

# 底部装饰线
add_shape(slide, Inches(1.5), Inches(5.5), Inches(3), Inches(0.06), ACCENT)


# ══════════════════════════════════════════════════════════════
# P2 · 为什么你的细胞需要鱼油？
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
            '为什么你的细胞需要鱼油？', font_size=44, bold=True, color=DARK)

# 分隔线
add_shape(slide, Inches(0.8), Inches(1.4), Inches(2), Inches(0.05), ACCENT)

# 左侧：问题
txBox = add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
                    '', font_size=20)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '现代人饮食的 Omega-3 危机'
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = ACCENT
p.font.name = 'Microsoft YaHei'

bullets = [
    '外卖、快餐大量使用大豆油、花生油 → Omega-6 严重超标',
    'Omega-6 与 Omega-3 理想比例 = 4:1 以下',
    '现代人实际比例高达 16:1 ~ 25:1',
    'Omega-6 过多 = 促炎（全身低度发炎）',
    '细胞膜被坏油糊满 → 营养进不去、垃圾出不来',
]
for b in bullets:
    add_paragraph(tf, f'  ▸  {b}', font_size=22, color=DARK, space_before=Pt(10))

# 右侧：概念框
box = add_shape(slide, Inches(7.2), Inches(2.0), Inches(5), Inches(4.5),
                RGBColor(0xEB, 0xF5, 0xFF))
txBox2 = add_textbox(slide, Inches(7.5), Inches(2.3), Inches(4.5), Inches(4.0),
                     '', font_size=20)
tf2 = txBox2.text_frame
tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = 'Omega-3 是什么？'
p2.font.size = Pt(28)
p2.font.bold = True
p2.font.color.rgb = ACCENT
p2.font.name = 'Microsoft YaHei'

bullets2 = [
    'EPA（二十碳五烯酸）→ 抗炎、血管清道夫',
    'DHA（二十二碳六烯酸）→ 脑神经、视网膜',
    '人体不能自行合成，必须从食物/补充剂获取',
    '补充 Omega-3 = 给细胞膜"换好油"',
]
for b in bullets2:
    add_paragraph(tf2, f'  ▸  {b}', font_size=22, color=DARK, space_before=Pt(10))


# ══════════════════════════════════════════════════════════════
# P3 · 鱼油的核心作用
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
            'Omega-3 的三大核心作用', font_size=44, bold=True, color=DARK)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(2), Inches(0.05), ACCENT)

# 三个卡片
cards = [
    ('🧹 换膜清洁', '置换细胞膜上的坏油\n恢复细胞膜流动性与弹性\n让营养进得去、垃圾出得来',
     ACCENT),
    ('🔥 全身抗炎', 'EPA 转化为消退素(Resolvins)\n直接关闭 NF-κB 炎症通路\n降低全身低度发炎水平',
     ACCENT2),
    ('🧠 神经保护', 'DHA 构成脑灰质与视网膜\n保护神经细胞髓鞘\n支持认知功能与视力',
     GREEN),
]

for i, (title, desc, color) in enumerate(cards):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2.0)
    # 卡片背景
    add_shape(slide, x, y, Inches(3.8), Inches(4.5), RGBColor(0xF8, 0xFA, 0xFC))
    # 顶部色条
    add_shape(slide, x, y, Inches(3.8), Inches(0.08), color)
    # 标题
    add_textbox(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.8),
                title, font_size=30, bold=True, color=color)
    # 描述
    txBox = add_textbox(slide, x + Inches(0.3), y + Inches(1.3), Inches(3.2), Inches(3.0),
                        '', font_size=20)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = desc.split('\n')
    for j, line in enumerate(lines):
        if j == 0:
            p = tf.paragraphs[0]
            p.text = line
        else:
            p = tf.add_paragraph()
            p.text = line
        p.font.size = Pt(22)
        p.font.color.rgb = DARK
        p.font.name = 'Microsoft YaHei'
        p.space_before = Pt(6)


# ══════════════════════════════════════════════════════════════
# P4 · 新旧鱼油对比
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(12), Inches(1),
            '新款高纯鱼油升级了什么？', font_size=44, bold=True, color=DARK)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(2), Inches(0.05), ACCENT)

# 对比表
headers = ['对比项', '原版深海鲑鱼油', '新款高纯鱼油', '升级亮点']
rows = [
    ['单粒 Omega-3', '300mg (EPA180+DHA120)', '345mg (更高纯度)', '↑ 15%'],
    ['每天建议量', '10粒 (3.0g)', '4粒 (1.38g)', '↓ 60% 粒数'],
    ['每天等效 Omega-3', '3.0g (治疗量)', '1.38g (基础抗炎)', '灵活搭配'],
    ['服用便利性', '需分3-4次服用', '随餐轻松4粒', '更方便'],
    ['搭配方案', '需单独大量补充', '组合1g+新款4粒=2.38g', '更灵活'],
]

# 表头
col_widths = [Inches(2.5), Inches(3.0), Inches(3.0), Inches(3.0)]
col_starts = [Inches(0.8)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

y_header = Inches(2.0)
add_shape(slide, Inches(0.8), y_header, Inches(11.5), Inches(0.8), ACCENT)
for i, (hdr, cs, cw) in enumerate(zip(headers, col_starts, col_widths)):
    add_textbox(slide, cs, y_header + Inches(0.05), cw, Inches(0.7),
                hdr, font_size=22, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# 数据行
for r_idx, row in enumerate(rows):
    y_row = y_header + Inches(0.8) + r_idx * Inches(0.85)
    bg_c = RGBColor(0xF8, 0xFA, 0xFC) if r_idx % 2 == 0 else WHITE
    add_shape(slide, Inches(0.8), y_row, Inches(11.5), Inches(0.85), bg_c)
    for c_idx, (cell, cs, cw) in enumerate(zip(row, col_starts, col_widths)):
        c = ACCENT if c_idx == 3 else DARK
        b = True if c_idx == 0 else False
        add_textbox(slide, cs, y_row + Inches(0.1), cw, Inches(0.65),
                    cell, font_size=20, bold=b, color=c, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# P5 · 产品规格
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
            '新款高纯鱼油 · 产品规格', font_size=44, bold=True, color=DARK)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(2), Inches(0.05), ACCENT)

# 核心数据：大号数字
data_items = [
    ('345', 'mg/粒', '高纯度\nOmega-3'),
    ('135', '粒/瓶', '每瓶含量'),
    ('4', '粒/天', '日常抗炎量'),
    ('1.38', 'g/天', '每日 Omega-3\n实际到账'),
]

for i, (num, unit, label) in enumerate(data_items):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.0)
    # 卡片
    add_shape(slide, x, y, Inches(2.8), Inches(3.5), RGBColor(0xF8, 0xFA, 0xFC))
    add_shape(slide, x, y, Inches(2.8), Inches(0.06), ACCENT)
    # 大数字
    add_textbox(slide, x, y + Inches(0.5), Inches(2.8), Inches(1.2),
                num, font_size=64, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    # 单位
    add_textbox(slide, x, y + Inches(1.6), Inches(2.8), Inches(0.6),
                unit, font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)
    # 标签
    add_textbox(slide, x + Inches(0.2), y + Inches(2.2), Inches(2.4), Inches(1.0),
                label, font_size=18, color=DARK, alignment=PP_ALIGN.CENTER)

# 底下说明
add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(1),
            '搭配男士活力组合（含1g Omega-3） → 全天 Omega-3 总量可达 2.38g，接近修复量 3.0g',
            font_size=22, color=GRAY)


# ══════════════════════════════════════════════════════════════
# P6 · 搭配方案
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
            '鱼油搭配方案 · 灵活选择', font_size=44, bold=True, color=DARK)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(2), Inches(0.05), ACCENT)

# 方案一：基础抗炎
y1 = Inches(2.0)
add_shape(slide, Inches(0.8), y1, Inches(5.5), Inches(4.5), RGBColor(0xEB, 0xF5, 0xFF))
add_shape(slide, Inches(0.8), y1, Inches(5.5), Inches(0.08), ACCENT)
add_textbox(slide, Inches(1.2), y1 + Inches(0.3), Inches(4.8), Inches(0.6),
            '方案一：日常抗炎保健', font_size=28, bold=True, color=ACCENT)
txBox = add_textbox(slide, Inches(1.2), y1 + Inches(1.1), Inches(4.8), Inches(3.0),
                    '', font_size=20)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '男士活力组合（2包）= 1g Omega-3'
p.font.size = Pt(22)
p.font.color.rgb = DARK
p.font.name = 'Microsoft YaHei'
p.space_before = Pt(6)
for line in ['+ 新款高纯鱼油 4粒 = 1.38g Omega-3', '', '合计：2.38g Omega-3/天', '适合：日常保健、基础抗炎']:
    add_paragraph(tf, line, font_size=22, color=DARK if '合计' not in line else ACCENT,
                  bold=True if '合计' in line else False, space_before=Pt(6))

# 方案二：治疗量
y2 = Inches(2.0)
add_shape(slide, Inches(7.0), y2, Inches(5.5), Inches(4.5), RGBColor(0xFD, 0xF0, 0xE0))
add_shape(slide, Inches(7.0), y2, Inches(5.5), Inches(0.08), ACCENT2)
add_textbox(slide, Inches(7.4), y2 + Inches(0.3), Inches(4.8), Inches(0.6),
            '方案二：细胞修复治疗量', font_size=28, bold=True, color=ACCENT2)
txBox = add_textbox(slide, Inches(7.4), y2 + Inches(1.1), Inches(4.8), Inches(3.0),
                    '', font_size=20)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '男士活力组合（2包）= 1g Omega-3'
p.font.size = Pt(22)
p.font.color.rgb = DARK
p.font.name = 'Microsoft YaHei'
p.space_before = Pt(6)
for line in ['+ 新款高纯鱼油 8粒 = 2.76g Omega-3', '', '合计：3.76g Omega-3/天', '适合：慢病修复、深度抗炎']:
    add_paragraph(tf, line, font_size=22, color=DARK if '合计' not in line else ACCENT2,
                  bold=True if '合计' in line else False, space_before=Pt(6))

# 底部说明
add_textbox(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.6),
            '细胞修复治疗量标准：Omega-3 3.0~4.0g/天 （石汉平教授《临床营养学》）',
            font_size=20, color=GRAY)


# ══════════════════════════════════════════════════════════════
# P7 · 经济账
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(1),
            '鱼油的经济账 · 值不值？', font_size=44, bold=True, color=DARK)
add_shape(slide, Inches(0.8), Inches(1.4), Inches(2), Inches(0.05), ACCENT)

# 定价卡片
# 左：产品定价
add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), RGBColor(0xEB, 0xF5, 0xFF))
add_shape(slide, Inches(0.8), Inches(2.0), Inches(0.08), Inches(4.5), ACCENT)

txBox = add_textbox(slide, Inches(1.2), Inches(2.3), Inches(4.8), Inches(4.0),
                    '', font_size=20)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '💰 定价'
p.font.size = Pt(30)
p.font.bold = True
p.font.color.rgb = ACCENT
p.font.name = 'Microsoft YaHei'

pricing = [
    '零售价：430 元 / 瓶',
    '每瓶 135 粒',
    '',
    '📊 每日成本',
    '每天 4 粒 → 12.7 元 / 天',
    '每天 8 粒 → 25.5 元 / 天',
    '',
    '📦 月花费',
    '4粒/天 → ≈383 元 / 月',
    '8粒/天 → ≈765 元 / 月',
]
for line in pricing:
    add_paragraph(tf, line, font_size=22, color=DARK, space_before=Pt(6))

# 右：对比
add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5), RGBColor(0xFD, 0xF0, 0xE0))
add_shape(slide, Inches(7.0), Inches(2.0), Inches(0.08), Inches(4.5), ACCENT2)

txBox = add_textbox(slide, Inches(7.4), Inches(2.3), Inches(4.8), Inches(4.0),
                    '', font_size=20)
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '⚖️ 对比：你每天都在花什么？'
p.font.size = Pt(26)
p.font.bold = True
p.font.color.rgb = ACCENT2
p.font.name = 'Microsoft YaHei'

comparisons = [
    '☕ 一杯奶茶 15 元 → 450 元 / 月',
    '🚬 一包烟 25 元 → 750 元 / 月',
    '🍔 一顿外卖 35 元 → 1,050 元 / 月',
    '',
    '✅ 鱼油 4粒/天 = 12.7 元 / 天',
    '   比一杯奶茶还便宜！',
    '',
    '🔑 本质：这不是消费，是投资',
    '   投资你的细胞健康',
]
for line in comparisons:
    add_paragraph(tf, line, font_size=22,
                  color=GRAY if '比一杯' in line else (ACCENT2 if '不是消费' in line else DARK),
                  bold=True if '不是消费' in line or '比一杯' in line else False,
                  space_before=Pt(6))


# ══════════════════════════════════════════════════════════════
# P8 · 总结 + 行动号召
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_shape(slide, Inches(0), Inches(0), Inches(0.3), H, ACCENT)

add_textbox(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(1.2),
            '总结：为什么选择新款高纯鱼油？', font_size=44, bold=True, color=WHITE)
add_shape(slide, Inches(1.5), Inches(2.1), Inches(2), Inches(0.05), ACCENT)

# 总结要点
points = [
    '✅  纯度更高 — 单粒 345mg Omega-3，超越原版',
    '✅  服用更方便 — 每天 4 粒即可达到基础抗炎量',
    '✅  方案灵活 — 可搭配活力组合，轻松升级到修复量',
    '✅  经济实惠 — 每天一杯奶茶钱，给全身细胞洗个澡',
    '✅  细胞级修复 — 换膜 + 抗炎 + 护神经，一油三用',
]

txBox = add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(3.5),
                    '', font_size=20)
tf = txBox.text_frame
tf.word_wrap = True
for i, point in enumerate(points):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = point
    p.font.size = Pt(26)
    p.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
    p.font.name = 'Microsoft YaHei'
    p.space_before = Pt(14)

# 行动号召
add_shape(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1.2), ACCENT)
add_textbox(slide, Inches(2.0), Inches(5.9), Inches(9), Inches(1.0),
            '想了解你适合哪个方案？私信我，帮你算一笔细胞修复的账',
            font_size=26, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ── 保存 ──
OUT_DIR = '/Users/mac/Documents/知识库/06-输出（成品）/ppt营养修复方案大众版'
os.makedirs(OUT_DIR, exist_ok=True)
out_path = os.path.join(OUT_DIR, '更新的鱼油-高纯Omega-3方案.pptx')
prs.save(out_path)
print(f'✅ PPT 已生成：{out_path}')
print(f'   共 {len(prs.slides)} 页')
