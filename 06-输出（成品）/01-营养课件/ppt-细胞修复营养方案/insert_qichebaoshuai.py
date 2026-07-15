#!/usr/bin/env python3
"""Insert 弃车保帅 page into 专业版 PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Open existing ──
src_path = '/Users/mac/Documents/知识库/06-输出（成品）/ppt-细胞修复营养方案/专业版细胞修复营养方案.pptx'
prs = Presentation(src_path)

# ── Colors ──
ACCENT  = RGBColor(0x00, 0x6D, 0xAA)
ACCENT2 = RGBColor(0xF0, 0x8C, 0x00)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x2E)
GRAY    = RGBColor(0x66, 0x66, 0x66)
GREEN   = RGBColor(0x27, 0xAE, 0x60)
RED     = RGBColor(0xE7, 0x4C, 0x3C)
CARD_BG = RGBColor(0xF8, 0xFA, 0xFC)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape

def add_rr(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, l, t, w, h, text, sz=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Microsoft YaHei'
    p.alignment = align
    return tb

# ── Create new slide: 弃车保帅 ──
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(sl, WHITE)

# Title
add_text(sl, Inches(0.6), Inches(0.3), Inches(11), Inches(0.7),
         '营养不够时，身体会"弃车保帅"', sz=32, bold=True, color=DARK)
add_shape(sl, Inches(0.6), Inches(0.95), Inches(1.5), Inches(0.04), ACCENT)
add_text(sl, Inches(0.6), Inches(1.1), Inches(11), Inches(0.4),
         'Bruce Ames Triage Theory — 细胞分流理论, PNAS 2006', sz=14, color=GRAY)

# Left panel: 生命器官
add_rr(sl, Inches(0.6), Inches(1.6), Inches(5.8), Inches(4.5), RGBColor(0xE8, 0xF8, 0xF0))
add_shape(sl, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.06), GREEN)
add_text(sl, Inches(1.0), Inches(1.8), Inches(5), Inches(0.5),
         '✅ 生命器官（优先供应）', sz=24, bold=True, color=GREEN)

vital = '❤️ 心脏    🫁 肺    🧠 大脑\n🫘 肝脏    🧊 肾脏'
add_text(sl, Inches(1.0), Inches(2.5), Inches(5), Inches(0.8),
         vital, sz=20, color=DARK)
add_text(sl, Inches(1.0), Inches(3.3), Inches(5), Inches(0.6),
         '少了任何一个你当场就活不了\n身体优先把营养分给它们', sz=16, color=GRAY)

# Right panel: 非生命器官
add_rr(sl, Inches(7.0), Inches(1.6), Inches(5.8), Inches(4.5), RGBColor(0xFD, 0xE0, 0xE0))
add_shape(sl, Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.06), RED)
add_text(sl, Inches(7.4), Inches(1.8), Inches(5), Inches(0.5),
         '⚠️ 非生命器官（被牺牲）', sz=24, bold=True, color=RED)

non_vital = '🫃 胃     🫀 乳腺    🦋 甲状腺\n🧬 子宫/卵巢  🦵 四肢    🧴 皮肤'
add_text(sl, Inches(7.4), Inches(2.5), Inches(5), Inches(0.8),
         non_vital, sz=20, color=DARK)
add_text(sl, Inches(7.4), Inches(3.3), Inches(5), Inches(0.6),
         '缺了它们不会马上死，但生活质量大幅下降\n营养不够，先牺牲它们', sz=16, color=GRAY)

# Symptoms strip at bottom
add_rr(sl, Inches(0.6), Inches(6.3), Inches(12), Inches(1.0), CARD_BG)
add_shape(sl, Inches(0.6), Inches(6.3), Inches(12), Inches(0.04), ACCENT2)
symp_text = (
    '反复口腔溃疡 → 胃粘膜缺原料  |  '
    '伤口愈合慢 → 缺蛋白/维C  |  '
    '皮肤松弛头发枯黄 → 缺胶原原料  |  '
    '乳腺增生/甲状腺结节反复 → 非生命器官拿不到营养'
)
add_text(sl, Inches(1.0), Inches(6.4), Inches(11), Inches(0.8),
         symp_text, sz=15, color=DARK)

# ═══ Reorder: move new slide to position 3 (after P2) ═══
sldIdLst = prs.slides._sldIdLst
sldId_elements = list(sldIdLst)
if len(sldId_elements) > 2:
    new_el = sldId_elements[-1]  # last element = newly added slide
    sldIdLst.remove(new_el)
    # insert after the 2nd element (index 1)
    ref_el = sldId_elements[1]
    ref_el.addnext(new_el)

out_path = '/Users/mac/Documents/知识库/06-输出（成品）/ppt-细胞修复营养方案/专业版细胞修复营养方案.pptx'
prs.save(out_path)
print(f'Done! Updated: {out_path}')
print(f'Total slides: {len(prs.slides)}')
